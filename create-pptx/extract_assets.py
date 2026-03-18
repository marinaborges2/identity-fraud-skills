#!/usr/bin/env python3
"""Extract Nubank logo assets from a reference PPTX presentation.

Usage:
    python extract_assets.py <path-to-nubank-presentation.pptx>

This extracts:
    - nu_logo.png       (small logo from cover slide, top-left)
    - section_logo.png  (larger logo from section divider slides)
    - cover_decoration.png (decorative image from cover slide, right side)
"""

import sys
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_assets.py <path-to-pptx>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    out_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
    os.makedirs(out_dir, exist_ok=True)

    prs = Presentation(pptx_path)

    # Slide 1: extract logo (small image) and cover decoration (large image)
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            w_in = shape.width / 914400
            blob = shape.image.blob
            ext = shape.image.content_type.split("/")[-1]
            if w_in < 1.5:
                path = os.path.join(out_dir, f"nu_logo.{ext}")
                with open(path, "wb") as f:
                    f.write(blob)
                print(f"Logo saved: {path} ({len(blob)} bytes)")
            else:
                path = os.path.join(out_dir, f"cover_decoration.{ext}")
                with open(path, "wb") as f:
                    f.write(blob)
                print(f"Cover decoration saved: {path} ({len(blob)} bytes)")

    # Look for section divider logo (slides with solid dark bg + logo)
    for slide_idx, slide in enumerate(prs.slides[1:], start=2):
        bg = slide.background
        if bg.fill.type is not None:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    w_in = shape.width / 914400
                    if 1.0 < w_in < 2.0:
                        blob = shape.image.blob
                        ext = shape.image.content_type.split("/")[-1]
                        path = os.path.join(out_dir, f"section_logo.{ext}")
                        with open(path, "wb") as f:
                            f.write(blob)
                        print(f"Section logo saved: {path} ({len(blob)} bytes)")
                        return

    print("\nDone! Assets saved to:", out_dir)


if __name__ == "__main__":
    main()
